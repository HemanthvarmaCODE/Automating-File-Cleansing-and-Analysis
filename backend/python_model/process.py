import os
import sys
import json
import re
import warnings
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed
import shutil
import tempfile

import pandas as pd
import fitz  
import docx
from pptx import Presentation
from openpyxl import load_workbook
import cv2
import pytesseract
from PIL import Image
import torch
from transformers import pipeline, BlipForConditionalGeneration, BlipProcessor

warnings.filterwarnings("ignore")

# -------------------------------
# Helpers for logging (stderr)
# -------------------------------
def log(msg: str):
    """Write informational/debug logs to stderr to avoid contaminating stdout JSON."""
    print(f"[INFO] {msg}", file=sys.stderr)

def error_log(msg: str):
    print(f"[ERROR] {msg}", file=sys.stderr)

# -------------------------------
# Configuration
# -------------------------------
MAX_WORKERS = min(8, (os.cpu_count() or 4)) 
CLEANSED_DIR = "cleansed_files"
REDACTED_IMG_DIR = "redacted_images"
OCR_CONF_THRESHOLD = 60  

os.makedirs(CLEANSED_DIR, exist_ok=True)
os.makedirs(REDACTED_IMG_DIR, exist_ok=True)

# -------------------------------
# Load models (GPU if available)
# -------------------------------
device_index = 0 if torch.cuda.is_available() else -1
log(f"torch.cuda.is_available: {torch.cuda.is_available()}, using device_index={device_index}")

try:
    ner_pipeline = pipeline(
        "ner",
        model="dslim/bert-base-NER",
        aggregation_strategy="simple",
        device=device_index,
    )
    blip_processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-large")
    blip_model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-large")
    if device_index >= 0:
        blip_model.to(f"cuda:{device_index}")
except Exception as e:
    error_log(f"Failed to load AI models: {e}")
    sys.stderr.write(json.dumps([{"error": f"Failed to load AI models: {e}"}]))
    sys.exit(1)

# -------------------------------
# Utility helpers
# -------------------------------
EMAIL_PATTERN = re.compile(r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b")

def safe_read_text_blocks_from_docx(doc):
    return [p.text for p in doc.paragraphs]

def safe_read_text_blocks_from_pptx(prs):
    blocks = []
    for s in prs.slides:
        for shape in s.shapes:
            if hasattr(shape, "text") and shape.text:
                blocks.append(shape.text)
    return blocks

def safe_read_text_blocks_from_xlsx(wb):
    blocks = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str):
                    blocks.append(cell.value)
    return blocks

def redact_using_map(orig_text, redaction_map):
    return redaction_map.get(orig_text, orig_text)

def analyze_and_redact_text_once(text: str):
    """
    Call NER once on the full `text` (a large string).
    Return redacted_text, pii_counts, vulnerabilities.
    """
    if not isinstance(text, str) or not text.strip():
        return text, {}, []

    try:
        entities = ner_pipeline(text) 
    except Exception as e:
        error_log(f"NER pipeline failure for a text block: {e}")
        return text, {}, [{"description": f"NER failed: {e}", "severity": "High"}]

    redacted = list(text)
    pii_counts = {"person": 0, "organization": 0, "location": 0, "email": 0, "phone": 0}
    vulnerabilities = []

    for ent in reversed(entities):
        start, end = ent.get("start"), ent.get("end")
        group = ent.get("entity_group")
        if group in ("PER", "ORG", "LOC") and start is not None and end is not None and start < end:
            redacted[start:end] = "[REDACTED]"
            if group == "PER":
                pii_counts["person"] += 1
            elif group == "ORG":
                pii_counts["organization"] += 1
            else:
                pii_counts["location"] += 1

    redacted_text = "".join(redacted)

    emails_found = len(EMAIL_PATTERN.findall(redacted_text))
    pii_counts["email"] = emails_found
    redacted_text = EMAIL_PATTERN.sub("[REDACTED]", redacted_text)

    if pii_counts["person"] > 0:
        vulnerabilities.append({"description": f"Detected {pii_counts['person']} person names.", "severity": "High"})
    if pii_counts["organization"] > 0:
        vulnerabilities.append({"description": f"Detected {pii_counts['organization']} organization names.", "severity": "Medium"})
    if pii_counts["email"] > 0:
        vulnerabilities.append({"description": f"Found {pii_counts['email']} email addresses.", "severity": "High"})

    return redacted_text, pii_counts, vulnerabilities

# -------------------------------
# File processors (optimized)
# -------------------------------
def process_pdf(file_path, output_path):
    try:
        doc = fitz.open(file_path)
    except Exception as e:
        error_log(f"Failed to open PDF {file_path}: {e}")
        return {}, [{"description": f"Failed to open PDF: {e}", "severity": "High"}]

    page_texts = [p.get_text() for p in doc]
    full_text = "\n".join([t for t in page_texts if t and t.strip()])

    redacted_full, pii_counts, vulnerabilities = analyze_and_redact_text_once(full_text)

    try:
        entities = ner_pipeline(full_text)
    except Exception:
        entities = []

    for page_idx, page in enumerate(doc):
        page_text = page_texts[page_idx] if page_idx < len(page_texts) else ""
        if not page_text or not page_text.strip():
            continue
        for ent in entities:
            if ent.get("entity_group") in ("PER", "ORG", "LOC"):
                word = ent.get("word")
                if not word:
                    continue
                try:
                    areas = page.search_for(word)
                    for rect in areas:
                        page.add_redact_annot(rect, fill=(0, 0, 0))
                except Exception:
                    continue
        try:
            page.apply_redactions()
        except Exception:
            continue

    try:
        doc.save(output_path, garbage=3, deflate=True)
    except Exception as e:
        error_log(f"Failed to save redacted PDF {output_path}: {e}")
        doc.close()
        return pii_counts, vulnerabilities + [{"description": f"Failed to save PDF: {e}", "severity": "High"}]
    doc.close()
    return pii_counts, vulnerabilities

def process_docx(file_path, output_path):
    try:
        doc = docx.Document(file_path)
    except Exception as e:
        error_log(f"Failed to open DOCX {file_path}: {e}")
        shutil.copy(file_path, output_path)
        return {}, [{"description": f"Failed to open DOCX: {e}", "severity": "High"}]

    blocks = safe_read_text_blocks_from_docx(doc)
    if not blocks:
        shutil.copy(file_path, output_path)
        return {}, []

    full_text = "\n".join(blocks)
    redacted_full, pii_counts, vulnerabilities = analyze_and_redact_text_once(full_text)
    redacted_lines = redacted_full.split("\n")

    redaction_map = {}
    for orig, red in zip(blocks, redacted_lines):
        if orig not in redaction_map:
            redaction_map[orig] = red

    changed = False
    for para in doc.paragraphs:
        if para.text and para.text.strip():
            new_text = redact_using_map(para.text, redaction_map)
            if new_text != para.text:
                para.text = new_text
                changed = True

    try:
        if changed:
            doc.save(output_path)
        else:
            shutil.copy(file_path, output_path)
    except Exception as e:
        error_log(f"Failed saving DOCX {output_path}: {e}")
        shutil.copy(file_path, output_path)
        return pii_counts, vulnerabilities + [{"description": f"Failed to save DOCX: {e}", "severity": "High"}]

    return pii_counts, vulnerabilities

def process_pptx(file_path, output_path):
    try:
        prs = Presentation(file_path)
    except Exception as e:
        error_log(f"Failed to open PPTX {file_path}: {e}")
        shutil.copy(file_path, output_path)
        return {}, [{"description": f"Failed to open PPTX: {e}", "severity": "High"}]

    blocks = safe_read_text_blocks_from_pptx(prs)
    if not blocks:
        shutil.copy(file_path, output_path)
        return {}, []

    full_text = "\n".join(blocks)
    redacted_full, pii_counts, vulnerabilities = analyze_and_redact_text_once(full_text)
    redacted_lines = redacted_full.split("\n")

    redaction_map = {}
    for orig, red in zip(blocks, redacted_lines):
        if orig not in redaction_map:
            redaction_map[orig] = red

    changed = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                new_text = redact_using_map(shape.text, redaction_map)
                if new_text != shape.text:
                    shape.text = new_text
                    changed = True

    try:
        if changed:
            prs.save(output_path)
        else:
            shutil.copy(file_path, output_path)
    except Exception as e:
        error_log(f"Failed saving PPTX {output_path}: {e}")
        shutil.copy(file_path, output_path)
        return pii_counts, vulnerabilities + [{"description": f"Failed to save PPTX: {e}", "severity": "High"}]

    return pii_counts, vulnerabilities

def process_xlsx(file_path, output_path):
    try:
        wb = load_workbook(file_path)
    except Exception as e:
        error_log(f"Failed to open XLSX {file_path}: {e}")
        shutil.copy(file_path, output_path)
        return {}, [{"description": f"Failed to open XLSX: {e}", "severity": "High"}]

    blocks = safe_read_text_blocks_from_xlsx(wb)
    if not blocks:
        shutil.copy(file_path, output_path)
        return {}, []

    full_text = "\n".join(blocks)
    redacted_full, pii_counts, vulnerabilities = analyze_and_redact_text_once(full_text)
    redacted_lines = redacted_full.split("\n")

    redaction_map = {}
    for orig, red in zip(blocks, redacted_lines):
        if orig not in redaction_map:
            redaction_map[orig] = red

    changed = False
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str):
                    new_text = redact_using_map(cell.value, redaction_map)
                    if new_text != cell.value:
                        cell.value = new_text
                        changed = True

    try:
        if changed:
            wb.save(output_path)
        else:
            shutil.copy(file_path, output_path)
    except Exception as e:
        error_log(f"Failed saving XLSX {output_path}: {e}")
        shutil.copy(file_path, output_path)
        return pii_counts, vulnerabilities + [{"description": f"Failed to save XLSX: {e}", "severity": "High"}]

    return pii_counts, vulnerabilities

def process_csv(file_path, output_path):
    try:
        df = pd.read_csv(file_path, dtype=str, keep_default_na=False)
    except Exception as e:
        error_log(f"Failed to read CSV {file_path}: {e}")
        shutil.copy(file_path, output_path)
        return {}, [{"description": f"Failed to read CSV: {e}", "severity": "High"}]

    row_blocks = df.astype(str).apply(lambda row: " ".join(row.values.tolist()), axis=1).tolist()
    if not any(r.strip() for r in row_blocks):
        shutil.copy(file_path, output_path)
        return {}, []

    full_text = "\n".join(row_blocks)
    redacted_full, pii_counts, vulnerabilities = analyze_and_redact_text_once(full_text)
    redacted_rows = redacted_full.split("\n")

    changed = False
    for col in df.columns:
        new_col = df[col].astype(str).apply(lambda v: analyze_and_redact_text_once(v)[0] if v and v.strip() else v)
        if not new_col.equals(df[col].astype(str)):
            df[col] = new_col
            changed = True

    try:
        if changed:
            df.to_csv(output_path, index=False)
        else:
            shutil.copy(file_path, output_path)
    except Exception as e:
        error_log(f"Failed saving CSV {output_path}: {e}")
        shutil.copy(file_path, output_path)
        return pii_counts, vulnerabilities + [{"description": f"Failed to save CSV: {e}", "severity": "High"}]

    return pii_counts, vulnerabilities

# -------------------------------
# Image processing (optimized)
# -------------------------------
def preprocess_for_ocr(cv_img):
    """Convert to grayscale + simple thresholding to improve OCR speed/accuracy."""
    gray = cv2.cvtColor(cv_img, cv2.COLOR_BGR2GRAY)
    _, th = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
    return th

def process_image(file_path, output_path):
    try:
        raw_image = Image.open(file_path).convert("RGB")
        inputs = blip_processor(raw_image, return_tensors="pt")
        input_device = next(blip_model.parameters()).device
        if getattr(input_device, "type", None) == "cuda":
            inputs = {k: v.to(input_device) for k, v in inputs.items()}
        out = blip_model.generate(**inputs)
        description = blip_processor.decode(out[0], skip_special_tokens=True)

        img = cv2.imread(file_path)
        if img is None:
            raise RuntimeError("cv2.imread failed (file corrupt or unreadable)")

        pre = preprocess_for_ocr(img)
        ocr_data = pytesseract.image_to_data(pre, output_type=pytesseract.Output.DICT)

        redacted_words = 0
        pii_counts = {}
        vulnerabilities = []
        ocr_text = " ".join([t for t in ocr_data.get("text", []) if t])
        redacted_set = set()
        if ocr_text.strip():
            ocr_redacted_full, ocr_pii_counts, ocr_vulns = analyze_and_redact_text_once(ocr_text)
            orig_tokens = ocr_text.split()
            red_tokens = ocr_redacted_full.split()
            for o, r in zip(orig_tokens, red_tokens):
                if o != r:
                    redacted_set.add(o)

        n = len(ocr_data.get("text", []))
        for i in range(n):
            try:
                conf = float(ocr_data["conf"][i])
            except Exception:
                conf = -1
            if conf < OCR_CONF_THRESHOLD:
                continue
            word = ocr_data["text"][i].strip()
            if not word:
                continue
            should_redact = (word in redacted_set) if redacted_set else False
            if not should_redact:
                red_word, _, _ = analyze_and_redact_text_once(word)
                should_redact = (red_word != word)
            if should_redact:
                redacted_words += 1
                (x, y, w, h) = (ocr_data["left"][i], ocr_data["top"][i], ocr_data["width"][i], ocr_data["height"][i])
                cv2.rectangle(img, (x, y), (x + w, y + h), (0, 0, 0), -1)

        if redacted_words > 0:
            vulnerabilities.append({"description": f"Redacted {redacted_words} sensitive words found via OCR.", "severity": "Medium"})

        try:
            if redacted_words > 0:
                cv2.imwrite(output_path, img)
            else:
                shutil.copy(file_path, output_path)
        except Exception as e:
            error_log(f"Failed saving image {output_path}: {e}")
            shutil.copy(file_path, output_path)
            return pii_counts, vulnerabilities + [{"description": f"Failed to save image: {e}", "severity": "High"}], description

        return pii_counts, vulnerabilities, description
    except Exception as e:
        error_log(f"Failed to process image {file_path}: {e}")
        return {}, [{"description": f"Failed to process image: {e}", "severity": "High"}], "Could not analyze image."

# -------------------------------
# Parallel orchestrator
# -------------------------------
def process_single_file_worker(directory_to_process, filename):
    file_path = os.path.join(directory_to_process, filename)
    if not os.path.isfile(file_path):
        return None
    p_file = Path(file_path)
    file_type = p_file.suffix.lower().replace(".", "")
    result = {
        "originalFileName": p_file.name,
        "fileType": file_type,
        "status": "completed",
        "summary": f"A .{file_type} file.",
        "piiDetected": {},
        "vulnerabilities": [],
        "cleansedFilePath": None,
    }
    try:
        cleansed_name = f"cleansed_{p_file.name}"
        if file_type in ["jpg", "jpeg", "png", "bmp"]:
            cleansed_path = os.path.join(REDACTED_IMG_DIR, cleansed_name)
        else:
            cleansed_path = os.path.join(CLEANSED_DIR, cleansed_name)

        if file_type == "pdf":
            pii, vulns = process_pdf(file_path, cleansed_path)
            result["summary"] = "PDF document analyzed and redacted."
        elif file_type == "docx":
            pii, vulns = process_docx(file_path, cleansed_path)
            result["summary"] = "Word document analyzed and redacted."
        elif file_type == "pptx":
            pii, vulns = process_pptx(file_path, cleansed_path)
            result["summary"] = "PowerPoint presentation analyzed and redacted."
        elif file_type in ("xlsx", "xlsm"):
            pii, vulns = process_xlsx(file_path, cleansed_path)
            result["summary"] = "Excel workbook analyzed and redacted."
        elif file_type == "csv":
            pii, vulns = process_csv(file_path, cleansed_path)
            result["summary"] = "CSV file analyzed and redacted."
        elif file_type in ["jpg", "jpeg", "png", "bmp"]:
            pii, vulns, description = process_image(file_path, cleansed_path)
            result["summary"] = description
        else:
            shutil.copy(file_path, cleansed_path)
            pii, vulns = {}, [{"description": f"File type '{file_type}' is not supported for redaction. File copied.", "severity": "Low"}]
            result["summary"] = "File format not supported for deep analysis."

        result["piiDetected"] = pii
        result["vulnerabilities"] = vulns
        result["cleansedFilePath"] = cleansed_path
    except Exception as e:
        error_log(f"Unexpected error processing {file_path}: {e}")
        result["status"] = "error"
        result["summary"] = "An error occurred during processing."
        result["vulnerabilities"] = [{"description": f"Critical error: {str(e)}", "severity": "High"}]

    return result

def main_parallel(directory_to_process, max_workers=MAX_WORKERS):
    all_results = []
    files = [f for f in os.listdir(directory_to_process) if os.path.isfile(os.path.join(directory_to_process, f))]
    if not files:
        return []

    log(f"Starting parallel processing of {len(files)} files with max_workers={max_workers}")
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        future_to_file = {executor.submit(process_single_file_worker, directory_to_process, f): f for f in files}
        for future in as_completed(future_to_file):
            try:
                res = future.result()
                if res is not None:
                    all_results.append(res)
            except Exception as e:
                error_log(f"Worker raised exception: {e}")
    log("Parallel processing complete")
    return all_results

# -------------------------------
# CLI entrypoint
# -------------------------------
if __name__ == "__main__":
    if len(sys.argv) < 2:
        error_log("Usage: python process.py <directory_to_process>")
        sys.exit(1)
    input_dir = sys.argv[1]
    if not os.path.isdir(input_dir):
        error_log(f"Directory not found: {input_dir}")
        sys.exit(1)

    results = main_parallel(input_dir)
    sys.stdout.write(json.dumps(results, indent=2))
    sys.stdout.flush()
